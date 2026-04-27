"""Generate chaingammon.pptx — ETHGlobal pitch deck."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x0D, 0x1A)   # near-black navy
ACCENT  = RGBColor(0x6C, 0x63, 0xFF)   # indigo/violet
ACCENT2 = RGBColor(0x00, 0xD4, 0xAA)   # teal
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GREY    = RGBColor(0xAA, 0xAA, 0xBB)
GOLD    = RGBColor(0xFF, 0xC8, 0x4B)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color=BG):
    """Fill slide background."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, left, top, width, height,
        fill=None, line=None, line_width=Pt(1)):
    """Add a coloured rectangle."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def txt(slide, text, left, top, width, height,
        size=Pt(18), bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    """Add a text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = size
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def bullet_frame(slide, items, left, top, width, height,
                 size=Pt(18), color=WHITE, marker="▸ ", spacing=1.2,
                 bold_first=False):
    """Add a multi-bullet text box."""
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = marker + item
        run.font.size  = size
        run.font.color.rgb = color
        run.font.bold  = (bold_first and i == 0)
    return txb


def accent_bar(slide, y=Inches(0.55), h=Inches(0.05)):
    """Thin accent line across the top."""
    bar = slide.shapes.add_shape(1, 0, y, W, h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def section_label(slide, text, y=Inches(0.12)):
    txt(slide, text, Inches(0.5), y, Inches(12), Inches(0.4),
        size=Pt(11), color=ACCENT, bold=True)


# ── slides ────────────────────────────────────────────────────────────────────

def slide_title(prs):
    """01 — Title"""
    s = blank_slide(prs)
    bg(s)

    # large diagonal accent block
    shape = s.shapes.add_shape(1, Inches(8.5), Inches(-0.5), Inches(5.5), Inches(9))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x16, 0x14, 0x30)
    shape.line.fill.background()

    accent_bar(s, y=Inches(0), h=Inches(0.5))

    txt(s, "CHAINGAMMON", Inches(0.6), Inches(1.2), Inches(8), Inches(1.2),
        size=Pt(54), bold=True, color=WHITE)

    txt(s, "An open protocol for portable backgammon reputation",
        Inches(0.6), Inches(2.5), Inches(8), Inches(0.8),
        size=Pt(22), color=ACCENT2)

    txt(s, "ETHGlobal Open Agents  ·  April 24 – May 6 2026",
        Inches(0.6), Inches(3.3), Inches(8), Inches(0.5),
        size=Pt(14), color=GREY, italic=True)

    # sponsor pills
    pills = [("ENS", ACCENT), ("0G Chain", ACCENT2), ("KeeperHub", GOLD)]
    x = Inches(0.6)
    for label, col in pills:
        b = s.shapes.add_shape(1, x, Inches(4.5), Inches(1.6), Inches(0.45))
        b.fill.solid(); b.fill.fore_color.rgb = col
        b.line.fill.background()
        tf = b.text_frame; tf.word_wrap = False
        p  = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r  = p.add_run(); r.text = label
        r.font.size = Pt(13); r.font.bold = True
        r.font.color.rgb = BG
        x += Inches(1.8)


def slide_problem(prs):
    """02 — The Problem"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "THE PROBLEM")

    txt(s, "Every backgammon platform is a walled graveyard.",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.75),
        size=Pt(34), bold=True, color=WHITE)

    # left: graveyard timeline
    txt(s, "Platform shutdowns — real losses, real players",
        Inches(0.5), Inches(1.75), Inches(6.0), Inches(0.4),
        size=Pt(14), bold=True, color=RGBColor(0xFF, 0x70, 0x70))

    events = [
        ("2014", "GameSpy", "Shut down Apr 30, 2014. All leaderboards, rankings, and statistics deleted overnight. No migration."),
        ("2020", "Microsoft Internet\nGames (Backgammon)", "Shut down Jan 22, 2020. 20 years of backgammon match history — permanently gone."),
        ("2026", "Facebook Gaming", "Creator Program ended Oct 2025; platform fully shutting down. Backgammon games and player data lost."),
    ]
    COL_A = RGBColor(0x1E, 0x10, 0x10)
    for i, (year, name, desc) in enumerate(events):
        y = Inches(2.25 + i * 1.5)
        row = s.shapes.add_shape(1, Inches(0.5), y, Inches(6.0), Inches(1.35))
        row.fill.solid(); row.fill.fore_color.rgb = COL_A; row.line.fill.background()
        pip = s.shapes.add_shape(1, Inches(0.5), y, Inches(0.22), Inches(1.35))
        pip.fill.solid(); pip.fill.fore_color.rgb = RGBColor(0xFF, 0x70, 0x70); pip.line.fill.background()
        txt(s, year, Inches(0.8), y + Inches(0.05), Inches(0.9), Inches(0.45),
            size=Pt(17), bold=True, color=RGBColor(0xFF, 0x70, 0x70))
        txt(s, name, Inches(1.75), y + Inches(0.05), Inches(4.6), Inches(0.45),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, desc, Inches(0.8), y + Inches(0.58), Inches(5.5), Inches(0.7),
            size=Pt(11), color=GREY)

    # right: fragmentation problem
    txt(s, "No standard. No portability. Nowhere safe.",
        Inches(7.0), Inches(1.75), Inches(5.8), Inches(0.4),
        size=Pt(14), bold=True, color=GOLD)

    frags = [
        "5+ active platforms today: Backgammon Galaxy, PlayOK, VIP Backgammon, Nextgammon, FIBS — ratings don\u2019t transfer between any of them.",
        "Chess has FIDE — a global, portable, federated rating body. Backgammon has nothing equivalent.",
        "Switch platforms \u2192 start at 1500. Decades of match history: reset. Your skill is real; your record is fiction.",
        "AI agents are even worse: no identity, no history, no verifiable track record across deployments.",
    ]
    bullet_frame(s, frags, Inches(7.0), Inches(2.25), Inches(5.9), Inches(3.3),
                 size=Pt(13), color=WHITE)

    # bottom quote bar
    q = s.shapes.add_shape(1, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.75))
    q.fill.solid(); q.fill.fore_color.rgb = RGBColor(0x1E, 0x1A, 0x42)
    q.line.color.rgb = ACCENT; q.line.width = Pt(1.5)
    tf = q.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = (
        '"California AB 2426 (2024) now requires 30 days\u2019 notice before a game shuts down.'
        ' It does NOT require preserving your ratings."'
    )
    r.font.size = Pt(13); r.font.italic = True; r.font.color.rgb = ACCENT2


def slide_market(prs):
    """03 — Market Opportunity"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "MARKET OPPORTUNITY")

    txt(s, "A Polymarket moment for skill — starting with backgammon",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(30), bold=True, color=WHITE)

    # ── left: Polymarket parallel ─────────────────────────────────────────────
    txt(s, "The Polymarket parallel", Inches(0.5), Inches(1.7), Inches(6.0), Inches(0.45),
        size=Pt(16), bold=True, color=ACCENT)

    pm_rows = [
        ("Polymarket",      "Chaingammon"),
        ("Centralised prediction\nmarkets (PredictIt)", "Centralised game ratings\n(Backgammon Galaxy, Chess.com)"),
        ("Opaque odds set\nby operators",                "ELO computed on-chain,\nverifiable by anyone"),
        ("Permissioned who can\ncreate a market",        "Permissionless: any platform\ncan read ENS text records"),
        ("Real $ at stake\naligns accuracy",             "Real match history on-chain\naligns trust"),
        ("Prices survive the\nplatform shutting down",   "Ratings survive the\nplatform shutting down"),
    ]
    COL_A = RGBColor(0x14, 0x12, 0x28)
    COL_B = RGBColor(0x11, 0x0F, 0x22)
    # header row
    hdr = s.shapes.add_shape(1, Inches(0.5), Inches(2.2), Inches(6.1), Inches(0.42))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = ACCENT; hdr.line.fill.background()
    txt(s, "BEFORE (Web2)", Inches(0.65), Inches(2.23), Inches(2.8), Inches(0.35),
        size=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)
    txt(s, "AFTER (Chaingammon)", Inches(3.55), Inches(2.23), Inches(2.9), Inches(0.35),
        size=Pt(11), bold=True, color=BG, align=PP_ALIGN.CENTER)

    for i, (before, after) in enumerate(pm_rows[1:]):   # skip header tuple
        y = Inches(2.67 + i * 0.74)
        row = s.shapes.add_shape(1, Inches(0.5), y, Inches(6.1), Inches(0.68))
        row.fill.solid()
        row.fill.fore_color.rgb = COL_A if i % 2 == 0 else COL_B
        row.line.fill.background()
        txt(s, before, Inches(0.65), y + Inches(0.04), Inches(2.75), Inches(0.6),
            size=Pt(11), color=GREY)
        txt(s, "\u2192", Inches(3.3), y + Inches(0.12), Inches(0.35), Inches(0.4),
            size=Pt(14), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        txt(s, after, Inches(3.7), y + Inches(0.04), Inches(2.75), Inches(0.6),
            size=Pt(11), bold=True, color=WHITE)

    # ── right: market size + extensibility ───────────────────────────────────
    txt(s, "Why backgammon first", Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.45),
        size=Pt(16), bold=True, color=ACCENT2)

    why_bg = [
        "Oldest recorded board game — 5,000+ year history, global player base",
        "ELO is well-understood and already used on every major platform",
        "gnubg: open-source world-class AI engine — a ready benchmark",
        "Skill gap between humans and AI is measurable, auditable, and meaningful",
        "Small, well-defined state space: every position is reproducible from dice + moves",
    ]
    bullet_frame(s, why_bg, Inches(7.0), Inches(2.2), Inches(5.9), Inches(2.8),
                 size=Pt(13), color=WHITE)

    txt(s, "Extensible to any turn-based game", Inches(7.0), Inches(5.1), Inches(5.8), Inches(0.45),
        size=Pt(15), bold=True, color=GOLD)

    games = [
        ("\u265f  Chess",    "FIDE ratings are already semi-portable; on-chain is the full solution"),
        ("\u25a1  Go",       "No global digital rating standard — same problem, larger community"),
        ("\u9006  Mahjong",  "Enormous Asian market; zero cross-platform rating portability"),
        ("\u2736  Any game", "Any game with reproducible positions + verifiable outcomes can use the protocol"),
    ]
    for i, (game, note) in enumerate(games):
        y = Inches(5.6 + i * 0.47)
        txt(s, game, Inches(7.0), y, Inches(1.4), Inches(0.42),
            size=Pt(13), bold=True, color=GOLD)
        txt(s, note, Inches(8.5), y, Inches(4.4), Inches(0.42),
            size=Pt(12), color=GREY)


def slide_solution(prs):
    """04 — The Solution"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "THE SOLUTION")

    txt(s, "Chaingammon — the open backgammon protocol",
        Inches(0.6), Inches(0.9), Inches(12), Inches(0.8),
        size=Pt(32), bold=True, color=WHITE)

    # three column cards
    cols = [
        ("🪪  Identity", "ENS subname\nalice.chaingammon.eth\ncarries ELO, match count,\nstyle URI, archive URI", ACCENT),
        ("🧠  Agent iNFT", "ERC-7857 iNFT on 0G Chain.\nEncrypted gnubg weights +\nlearned experience overlay,\nboth on 0G Storage.", ACCENT2),
        ("⚖️  Settlement", "KeeperHub workflow:\nrecordMatch + ENS text\nrecords + overlay update +\nverifiable audit trail.", GOLD),
    ]
    for i, (title, body, col) in enumerate(cols):
        x = Inches(0.5 + i * 4.25)
        card = s.shapes.add_shape(1, x, Inches(2.1), Inches(4.0), Inches(4.5))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
        card.line.color.rgb = col; card.line.width = Pt(1.5)

        top_bar = s.shapes.add_shape(1, x, Inches(2.1), Inches(4.0), Inches(0.35))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = col
        top_bar.line.fill.background()

        txt(s, title, x + Inches(0.15), Inches(2.55), Inches(3.7), Inches(0.6),
            size=Pt(16), bold=True, color=WHITE)
        txt(s, body, x + Inches(0.15), Inches(3.2), Inches(3.7), Inches(2.8),
            size=Pt(14), color=GREY)


def slide_how_it_works(prs):
    """04 — How It Works"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "HOW IT WORKS")

    txt(s, "End-to-end match lifecycle",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)

    steps = [
        ("1", "Connect wallet", "Frontend auto-issues alice.chaingammon.eth ENS subname"),
        ("2", "Pick opponent", "Another human subname, or an AI agent iNFT (gnubg-advanced-1.chaingammon.eth)"),
        ("3", "Play", "Server rolls dice; gnubg weights pulled from 0G Storage, run server-side per turn"),
        ("4", "Game ends", "Server builds GameRecord JSON → uploads to 0G Storage → gets Merkle rootHash"),
        ("5", "KeeperHub fires", "recordMatch on-chain (with rootHash) + ENS text records updated + agent overlay refreshed"),
        ("6", "Audit published", "KeeperHub audit JSON mirrored to 0G Storage; match replay publicly accessible forever"),
    ]
    for i, (num, title, desc) in enumerate(steps):
        row = i // 3
        col = i % 3
        x = Inches(0.5 + col * 4.25)
        y = Inches(1.8 + row * 2.7)

        circ = s.shapes.add_shape(9, x, y, Inches(0.45), Inches(0.45))   # oval
        circ.fill.solid(); circ.fill.fore_color.rgb = ACCENT
        circ.line.fill.background()
        tf = circ.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = num; r.font.bold = True; r.font.size = Pt(13)
        r.font.color.rgb = WHITE

        txt(s, title,  x + Inches(0.55), y, Inches(3.5), Inches(0.45),
            size=Pt(15), bold=True, color=WHITE)
        txt(s, desc, x + Inches(0.55), y + Inches(0.45), Inches(3.5), Inches(1.8),
            size=Pt(12), color=GREY)


def slide_ens(prs):
    """05 — ENS & Subnames"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "SPONSOR SPOTLIGHT — ENS")

    txt(s, "ENS subnames as portable player identity",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(30), bold=True, color=WHITE)

    # ── left column: what ENS is ──────────────────────────────────────────────
    txt(s, "What is ENS?", Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.45),
        size=Pt(16), bold=True, color=ACCENT)

    ens_intro = [
        "Ethereum Name Service — the DNS of Web3",
        "Human-readable names (alice.eth) resolve to addresses, content hashes, and arbitrary text records",
        "Subnames: any .eth owner can issue <label>.<name>.eth — delegated, on-chain, free to transfer",
        "Text records are on-chain structured metadata: a player is just a name + a bag of key-value facts",
    ]
    bullet_frame(s, ens_intro, Inches(0.5), Inches(2.2), Inches(5.8), Inches(2.5),
                 size=Pt(13), color=WHITE)

    txt(s, "Why subnames are the right primitive",
        Inches(0.5), Inches(4.85), Inches(5.8), Inches(0.45),
        size=Pt(14), bold=True, color=ACCENT)
    why = [
        "Platform-portable: alice.chaingammon.eth exists even if Chaingammon.com goes dark",
        "Composable: any dApp reads text records permissionlessly via ENS resolver",
        "Human-readable: beats 0x4a…3f in leaderboards, match UIs, and social sharing",
        "Hierarchical trust: chaingammon.eth controls issuance; player controls the name",
    ]
    bullet_frame(s, why, Inches(0.5), Inches(5.35), Inches(5.8), Inches(2.0),
                 size=Pt(12), color=GREY)

    # ── right column: our text records ───────────────────────────────────────
    txt(s, "alice.chaingammon.eth  text records",
        Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.45),
        size=Pt(16), bold=True, color=ACCENT2)

    records = [
        ("elo",            "1847",                  "Current ELO rating — updated after every match"),
        ("match_count",    "42",                    "Total games played, on-chain provenance"),
        ("last_match_id",  "0x1f…c3",               "Pointer to latest MatchRegistry entry"),
        ("style_uri",      "0g://kv/alice/style",   "Link to 0G KV style profile (opening tendencies, cube aggression)"),
        ("archive_uri",    "0g://log/alice/games",  "Link to full 0G Log of all game records"),
    ]
    COL_A = RGBColor(0x14, 0x12, 0x28)
    COL_B = RGBColor(0x11, 0x0F, 0x22)
    for i, (key, val, note) in enumerate(records):
        y = Inches(2.25 + i * 0.96)
        row = s.shapes.add_shape(1, Inches(7.0), y, Inches(5.9), Inches(0.88))
        row.fill.solid()
        row.fill.fore_color.rgb = COL_A if i % 2 == 0 else COL_B
        row.line.fill.background()
        txt(s, key, Inches(7.12), y + Inches(0.04), Inches(1.5), Inches(0.35),
            size=Pt(12), bold=True, color=ACCENT2)
        txt(s, val, Inches(8.75), y + Inches(0.04), Inches(2.05), Inches(0.35),
            size=Pt(12), color=GOLD)
        txt(s, note, Inches(7.12), y + Inches(0.44), Inches(5.65), Inches(0.38),
            size=Pt(10.5), color=GREY)

    txt(s, "PlayerSubnameRegistrar.sol issues names; KeeperHub writes text records post-match.",
        Inches(7.0), Inches(7.1), Inches(5.9), Inches(0.35),
        size=Pt(11), color=GREY, italic=True)


def slide_0g(prs):
    """06 — 0G Chain & Storage"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "SPONSOR SPOTLIGHT — 0G CHAIN & STORAGE")

    txt(s, "AI-native modular blockchain — chain + storage in one stack",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)

    # ── top row: two platform cards ───────────────────────────────────────────
    platform_cards = [
        ("0G Chain", ACCENT,
         "EVM-compatible L1  \u00b7  chainId 16602\n"
         "Fast finality, EVM Cancun opcodes.\n"
         "Hosts MatchRegistry, AgentRegistry, EloMath, PlayerSubnameRegistrar.\n"
         "Verified on chainscan-galileo block explorer."),
        ("0G Storage", ACCENT2,
         "Decentralized, content-addressed, permanent.\n"
         "Three primitives: Log (append-only), KV (key-value), Blob (binary).\n"
         "Upload returns a root hash you can commit on-chain.\n"
         "Designed for AI/ML workloads \u2014 ideal for model weights storage."),
    ]
    for i, (title, col, body) in enumerate(platform_cards):
        x = Inches(0.5 + i * 6.4)
        card = s.shapes.add_shape(1, x, Inches(1.75), Inches(6.1), Inches(2.3))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
        card.line.color.rgb = col; card.line.width = Pt(1.5)
        tag = s.shapes.add_shape(1, x, Inches(1.75), Inches(6.1), Inches(0.35))
        tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background()
        txt(s, title, x + Inches(0.15), Inches(2.18), Inches(5.8), Inches(0.45),
            size=Pt(17), bold=True, color=WHITE)
        txt(s, body, x + Inches(0.15), Inches(2.65), Inches(5.8), Inches(1.3),
            size=Pt(12), color=GREY)

    # ── bottom: why 0G fits ───────────────────────────────────────────────────
    txt(s, "Why 0G is the right fit for Chaingammon",
        Inches(0.5), Inches(4.2), Inches(12), Inches(0.45),
        size=Pt(16), bold=True, color=GOLD)

    why_rows = [
        ("Content-addressed linking",
         "rootHash returned by 0G Storage is committed in MatchRegistry.recordMatch — chain \u21d4 storage in one atomic reference."),
        ("AI-native storage",
         "gnubg neural-network weights (~408 KB) and per-agent experience overlays live in 0G Blob. 0G was built for exactly this."),
        ("Immutable audit trail",
         "Log type is append-only. Game records and KeeperHub audit JSONs persist forever without IPFS pinning risk."),
        ("One chain, one SDK",
         "og-bridge wraps @0gfoundation/0g-ts-sdk. Same RPC for contracts and storage \u2014 no cross-chain bridging ceremony."),
    ]
    COL_A = RGBColor(0x14, 0x12, 0x28)
    COL_B = RGBColor(0x11, 0x0F, 0x22)
    for i, (title, body) in enumerate(why_rows):
        col_idx = i % 2
        row_idx = i // 2
        x = Inches(0.5 + col_idx * 6.4)
        y = Inches(4.75 + row_idx * 1.1)
        row = s.shapes.add_shape(1, x, y, Inches(6.1), Inches(1.0))
        row.fill.solid()
        row.fill.fore_color.rgb = COL_A if i % 2 == 0 else COL_B
        row.line.fill.background()
        txt(s, title, x + Inches(0.15), y + Inches(0.05), Inches(5.8), Inches(0.35),
            size=Pt(13), bold=True, color=ACCENT2)
        txt(s, body, x + Inches(0.15), y + Inches(0.44), Inches(5.8), Inches(0.5),
            size=Pt(11), color=GREY)


def slide_keeperhub_intro(prs):
    """07 — KeeperHub"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "SPONSOR SPOTLIGHT — KEEPERHUB")

    txt(s, "Programmable keeper automation for multi-step Web3 workflows",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(26), bold=True, color=WHITE)

    # ── left: what is KeeperHub ───────────────────────────────────────────────
    txt(s, "What is KeeperHub?", Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.45),
        size=Pt(16), bold=True, color=GOLD)

    what_bullets = [
        "Keeper network: third-party bots watch for triggers and execute arbitrary on-chain transactions",
        "Workflow engine: define ordered steps with retry logic, gas configuration, and conditional branching",
        "Triggers: HTTP webhooks, on-chain events, cron schedules",
        "Audit-first: every workflow run produces a full execution trace with timestamps, retries, and gas used",
        "Decoupled from the app: KeeperHub will execute even if your server goes offline after firing the trigger",
    ]
    bullet_frame(s, what_bullets, Inches(0.5), Inches(2.2), Inches(5.9), Inches(3.2),
                 size=Pt(13), color=WHITE)

    txt(s, "vs. doing settlement in the server:",
        Inches(0.5), Inches(5.55), Inches(5.8), Inches(0.4),
        size=Pt(13), bold=True, color=ACCENT)
    compare = [
        "Server crashes mid-settlement \u2192 KeeperHub retries independently",
        "Gas spikes \u2192 KeeperHub optimizes and resubmits",
        "4 contract calls in sequence \u2192 one atomic workflow definition",
        "Compliance / audit \u2192 signed execution receipts, not server logs",
    ]
    bullet_frame(s, compare, Inches(0.5), Inches(6.0), Inches(5.9), Inches(1.4),
                 size=Pt(12), color=GREY)

    # ── right: our workflow ───────────────────────────────────────────────────
    txt(s, "Chaingammon settlement workflow",
        Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.45),
        size=Pt(16), bold=True, color=GOLD)

    TRIGGER_COL = RGBColor(0x14, 0x12, 0x28)
    steps_kh = [
        (ACCENT,   "HTTP trigger",        "Server calls POST /trigger after game ends"),
        (ACCENT,   "recordMatch",         "MatchRegistry.recordMatch(winner, loser, length, gameRecordHash) \u2192 emits EloUpdated"),
        (ACCENT2,  "ENS text records",    "setText on alice.chaingammon.eth and opponent: elo, match_count, last_match_id, archive_uri"),
        (GOLD,     "updateOverlayHash",   "AgentRegistry.updateOverlayHash(agentId, newHash) bumps experienceVersion"),
        (RGBColor(0xC0,0x80,0xFF), "Audit export", "Full trace \u2192 0G Storage Log \u2192 publicly viewable settlement proof"),
    ]
    for i, (col, title, body) in enumerate(steps_kh):
        y = Inches(2.2 + i * 0.97)
        row = s.shapes.add_shape(1, Inches(7.0), y, Inches(5.9), Inches(0.88))
        row.fill.solid(); row.fill.fore_color.rgb = TRIGGER_COL; row.line.fill.background()
        pip = s.shapes.add_shape(1, Inches(7.0), y, Inches(0.2), Inches(0.88))
        pip.fill.solid(); pip.fill.fore_color.rgb = col; pip.line.fill.background()
        txt(s, title, Inches(7.3), y + Inches(0.04), Inches(5.5), Inches(0.35),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, body, Inches(7.3), y + Inches(0.44), Inches(5.5), Inches(0.38),
            size=Pt(11), color=GREY)
        if i < len(steps_kh) - 1:
            txt(s, "\u2193", Inches(7.08), y + Inches(0.88), Inches(0.25), Inches(0.2),
                size=Pt(10), color=GREY, align=PP_ALIGN.CENTER)

    txt(s, "kh CLI defines the workflow YAML; the server only fires the trigger.",
        Inches(7.0), Inches(7.1), Inches(5.9), Inches(0.35),
        size=Pt(11), color=GREY, italic=True)


def slide_agents(prs):
    """08 — AI Agents as iNFTs"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "AI AGENTS AS iNFTs  (ERC-7857)")

    txt(s, "Agents that carry verifiable, portable intelligence",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)

    # left: two data hashes
    txt(s, "Two data hashes per iNFT", Inches(0.6), Inches(1.75), Inches(5.8), Inches(0.5),
        size=Pt(17), bold=True, color=ACCENT2)

    hash_rows = [
        ("dataHashes[0]", "baseWeightsHash", "Shared across ALL agents — encrypted gnubg\nneural-network weights on 0G Storage.\nVerifies: 'this agent runs real gnubg.'"),
        ("dataHashes[1]", "overlayHash", "Per-agent experience overlay — a ~20-float\nvector that biases move selection and grows\nwith every match played."),
    ]
    for i, (key, name, desc) in enumerate(hash_rows):
        y = Inches(2.35 + i * 2.1)
        card = s.shapes.add_shape(1, Inches(0.5), y, Inches(6.0), Inches(1.85))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
        card.line.color.rgb = ACCENT; card.line.width = Pt(1)
        txt(s, key,  Inches(0.7), y + Inches(0.1), Inches(5.5), Inches(0.35),
            size=Pt(11), color=ACCENT, bold=True)
        txt(s, name, Inches(0.7), y + Inches(0.4), Inches(5.5), Inches(0.35),
            size=Pt(14), bold=True, color=WHITE)
        txt(s, desc, Inches(0.7), y + Inches(0.75), Inches(5.5), Inches(1.0),
            size=Pt(12), color=GREY)

    # right: tier table
    txt(s, "Tier system", Inches(7.1), Inches(1.75), Inches(5.5), Inches(0.5),
        size=Pt(17), bold=True, color=GOLD)

    tiers = [
        ("Tier 0", "Beginner",      "gnubg 0-ply (immediate eval)"),
        ("Tier 1", "Intermediate",  "gnubg 2-ply (looks 2 moves ahead)"),
        ("Tier 2", "Advanced",      "gnubg 3-ply (deeper search)"),
        ("Tier 3", "World-class",   "gnubg full rollout"),
    ]
    for i, (t, label, note) in enumerate(tiers):
        y = Inches(2.35 + i * 1.1)
        row = s.shapes.add_shape(1, Inches(7.1), y, Inches(5.8), Inches(0.95))
        row.fill.solid()
        row.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28) if i % 2 == 0 else RGBColor(0x11, 0x0F, 0x22)
        row.line.fill.background()
        txt(s, t,     Inches(7.2), y + Inches(0.05), Inches(0.9),  Inches(0.4),
            size=Pt(13), bold=True, color=GOLD)
        txt(s, label, Inches(8.2), y + Inches(0.05), Inches(2.0),  Inches(0.4),
            size=Pt(13), bold=True, color=WHITE)
        txt(s, note,  Inches(7.2), y + Inches(0.45), Inches(4.5),  Inches(0.45),
            size=Pt(11), color=GREY)

    txt(s, "Two same-tier agents start identically — then drift into different styles as their match histories diverge.",
        Inches(7.1), Inches(6.65), Inches(5.8), Inches(0.7),
        size=Pt(12), color=ACCENT2, italic=True)


def slide_elo(prs):
    """06 — ELO Ratings & Prediction Markets"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "ELO RATINGS & PREDICTION MARKETS")

    txt(s, "A transparent, on-chain skill ledger — and a betting primitive",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(26), bold=True, color=WHITE)

    # ── LEFT: how ELO works ───────────────────────────────────────────────────
    txt(s, "How ELO works", Inches(0.5), Inches(1.7), Inches(5.8), Inches(0.5),
        size=Pt(17), bold=True, color=ACCENT2)

    # formula card
    form_card = s.shapes.add_shape(1, Inches(0.5), Inches(2.25), Inches(5.9), Inches(1.7))
    form_card.fill.solid(); form_card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
    form_card.line.color.rgb = ACCENT2; form_card.line.width = Pt(1.2)

    txt(s, "E\u2090 = 1 / (1 + 10^((R_B \u2212 R_A) / 400))",
        Inches(0.65), Inches(2.35), Inches(5.6), Inches(0.5),
        size=Pt(14), bold=True, color=ACCENT2)
    txt(s, "R_new = R_old + K \u00d7 (score \u2212 E\u2090)",
        Inches(0.65), Inches(2.85), Inches(5.6), Inches(0.45),
        size=Pt(14), bold=True, color=WHITE)
    txt(s, "K = 32 (rating change cap)   \u2022   Starting rating = 1500",
        Inches(0.65), Inches(3.3), Inches(5.6), Inches(0.4),
        size=Pt(12), color=GREY)

    # probability table
    txt(s, "Rating gap  \u2192  win probability",
        Inches(0.5), Inches(4.1), Inches(5.8), Inches(0.4),
        size=Pt(13), bold=True, color=ACCENT)

    prob_rows = [
        ("\u03940  (equal)",    "50%"),
        ("\u0394100",           "64%"),
        ("\u0394200",           "76%"),
        ("\u0394400",           "91%"),
    ]
    COL_A = RGBColor(0x14, 0x12, 0x28)
    COL_B = RGBColor(0x11, 0x0F, 0x22)
    for i, (delta, pct) in enumerate(prob_rows):
        y = Inches(4.55 + i * 0.48)
        row = s.shapes.add_shape(1, Inches(0.5), y, Inches(5.9), Inches(0.43))
        row.fill.solid()
        row.fill.fore_color.rgb = COL_A if i % 2 == 0 else COL_B
        row.line.fill.background()
        txt(s, delta, Inches(0.65), y + Inches(0.04), Inches(3.5), Inches(0.38),
            size=Pt(12), color=WHITE)
        txt(s, pct, Inches(4.5), y + Inches(0.04), Inches(1.6), Inches(0.38),
            size=Pt(12), bold=True, color=GOLD, align=PP_ALIGN.RIGHT)

    txt(s, "Chaingammon uses fixed-point integer ELO via EloMath.sol — no floats on-chain.",
        Inches(0.5), Inches(6.65), Inches(5.9), Inches(0.6),
        size=Pt(11), color=GREY, italic=True)

    # ── RIGHT: betting mechanisms ─────────────────────────────────────────────
    txt(s, "Betting on players & agents", Inches(7.0), Inches(1.7), Inches(5.8), Inches(0.5),
        size=Pt(17), bold=True, color=GOLD)

    bets = [
        ("Match odds", "Derive fair odds directly from the ELO win-probability formula."
         " A 200-point gap gives 76 / 24 fair odds — any market price away from that is +EV."),
        ("ELO movement markets", "Bet whether a player gains or loses N ELO over the next M games."
         " Think prediction markets on skill trajectory."),
        ("Equity-curve exposure", "Buy a share of an agent\u2019s future ELO gains."
         " Agents with long on-chain histories show clean cumulative curves"
         " \u2014 verifiable Sharpe for a bot."),
        ("Kelly sizing", "Optimal bet% = edge / odds. With ELO-implied probability and"
         " on-chain match history, the edge estimate is transparent, not operator-controlled."),
        ("Agent performance futures", "Will an agent reach Tier 3 ELO within 100 matches?"
         " All inputs\u2014match count, overlay hash\u2014are public on 0G Chain."),
    ]

    y_start = Inches(2.25)
    for i, (title, body) in enumerate(bets):
        y = y_start + Inches(i * 0.99)
        pill = s.shapes.add_shape(1, Inches(7.0), y, Inches(5.9), Inches(0.88))
        pill.fill.solid()
        pill.fill.fore_color.rgb = COL_A if i % 2 == 0 else COL_B
        pill.line.fill.background()
        txt(s, title, Inches(7.15), y + Inches(0.03), Inches(5.6), Inches(0.35),
            size=Pt(13), bold=True, color=GOLD)
        txt(s, body, Inches(7.15), y + Inches(0.4), Inches(5.6), Inches(0.45),
            size=Pt(10.5), color=GREY)

    txt(s, "Every match, ELO delta, and overlay update is recorded on 0G Chain"
        " \u2014 no trust required.",
        Inches(7.0), Inches(7.1), Inches(5.9), Inches(0.35),
        size=Pt(11), color=ACCENT2, italic=True)


def slide_contracts(prs):
    """07 — Smart Contracts on 0G Chain  (was 06)"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "SMART CONTRACTS ON 0G CHAIN  (chainId 16602)")

    txt(s, "Four contracts — deployed and verified on chainscan-galileo",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(26), bold=True, color=WHITE)

    contracts = [
        ("EloMath.sol", ACCENT,
         "Pure library. K=32, INITIAL=1500 ELO.\nFixed-point integer math — no floats on-chain.\nExpected-score lookup table, floor at 0."),
        ("MatchRegistry.sol", ACCENT2,
         "Owner-only recordMatch(winner, loser, matchLength, gameRecordHash).\nStores MatchInfo struct with on-chain ↔ 0G Storage link.\nEmits MatchRecorded, EloUpdated, GameRecordStored events."),
        ("AgentRegistry.sol", GOLD,
         "ERC-721 base + ERC-7857-compatible shape.\nmintAgent(to, metadataURI, tier)  ·  setBaseWeightsHash\nupdateOverlayHash — bumps matchCount + experienceVersion."),
        ("PlayerSubnameRegistrar.sol", RGBColor(0xFF, 0x70, 0x70),
         "Issues <name>.chaingammon.eth ENS subnames.\nServer-only write access in v1.\nControls text records: elo, match_count, style_uri, archive_uri."),
    ]

    for i, (name, col, desc) in enumerate(contracts):
        row = i // 2; c = i % 2
        x = Inches(0.5 + c * 6.4)
        y = Inches(1.8 + row * 2.6)
        card = s.shapes.add_shape(1, x, y, Inches(6.1), Inches(2.3))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
        card.line.color.rgb = col; card.line.width = Pt(1.5)
        tag = s.shapes.add_shape(1, x, y, Inches(6.1), Inches(0.35))
        tag.fill.solid(); tag.fill.fore_color.rgb = col; tag.line.fill.background()
        txt(s, name, x + Inches(0.12), y + Inches(0.38), Inches(5.8), Inches(0.45),
            size=Pt(15), bold=True, color=WHITE)
        txt(s, desc, x + Inches(0.12), y + Inches(0.85), Inches(5.8), Inches(1.4),
            size=Pt(12), color=GREY)


def slide_storage(prs):
    """07 — 0G Storage"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "0G STORAGE — FOUR USES")

    txt(s, "Decentralised, content-addressed, permanent",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)

    rows = [
        ("Log — per match",    "Game record",       "GameRecord JSON: every move, dice, cube action, gnubg position IDs. ~2–10 KB. rootHash committed in MatchRegistry on-chain.",           ACCENT),
        ("Log — per match",    "KeeperHub audit",   "Settlement workflow trace: timestamps, retries, gas, step results. Mirrored so it's publicly viewable without KeeperHub login.",         ACCENT2),
        ("KV — per player",    "Style profile",     "Descriptive stats only: % opening style, cube tendency, win rate vs each tier. Displayed on the player profile page.",                   GOLD),
        ("Blob — shared",      "gnubg base weights","Encrypted (AES-256-GCM). ~408 KB. One blob for all agents. rootHash is dataHashes[0] on every iNFT. Key held server-side in v1.",      RGBColor(0xFF, 0x70, 0x70)),
        ("Blob — per agent",   "Experience overlay","~20-float preference vector. Updates after every match. rootHash is dataHashes[1] on the iNFT. Grows with play.",                        RGBColor(0xC0, 0x80, 0xFF)),
    ]

    for i, (kind, name, desc, col) in enumerate(rows):
        y = Inches(1.75 + i * 1.05)
        stripe = s.shapes.add_shape(1, Inches(0.3), y, Inches(0.18), Inches(0.85))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = col; stripe.line.fill.background()
        txt(s, kind, Inches(0.6), y + Inches(0.02), Inches(2.0), Inches(0.4),
            size=Pt(11), color=col, bold=True)
        txt(s, name, Inches(2.7), y + Inches(0.02), Inches(2.5), Inches(0.4),
            size=Pt(14), bold=True, color=WHITE)
        txt(s, desc, Inches(5.4), y + Inches(0.02), Inches(7.5), Inches(0.85),
            size=Pt(12), color=GREY)
        box(slide=s, left=Inches(0.3), top=y + Inches(0.87),
            width=Inches(12.7), height=Pt(0.5),
            fill=RGBColor(0x1C, 0x19, 0x35))


def slide_keeperhub(prs):
    """08 — KeeperHub Settlement"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "KEEPERHUB — SETTLEMENT WORKFLOW")

    txt(s, "Multi-step orchestration with retry, gas optimization, and audit",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(25), bold=True, color=WHITE)

    # flow
    steps_kh = [
        ("1", "recordMatch", "MatchRegistry.recordMatch(\n  winner, loser, matchLength,\n  gameRecordHash\n)\nReturns matchId; ELO updated on-chain.", ACCENT),
        ("2", "ENS update", "setText(alice.chaingammon.eth,\n  elo, match_count,\n  last_match_id, archive_uri)\nfor both players.", ACCENT2),
        ("3", "Overlay update", "AgentRegistry.updateOverlayHash(\n  agentId, newOverlayHash\n)\nbumps experienceVersion +\nmatchCount on the iNFT.", GOLD),
        ("4", "Audit JSON", "Workflow emits full trace:\ntimestamps, retry count, gas\nper step. Mirrored to\n0G Storage alongside game record.", RGBColor(0xC0, 0x80, 0xFF)),
    ]

    for i, (num, title, body, col) in enumerate(steps_kh):
        x = Inches(0.5 + i * 3.2)
        card = s.shapes.add_shape(1, x, Inches(1.9), Inches(3.0), Inches(4.3))
        card.fill.solid(); card.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
        card.line.color.rgb = col; card.line.width = Pt(1.5)

        top_bar = s.shapes.add_shape(1, x, Inches(1.9), Inches(3.0), Inches(0.35))
        top_bar.fill.solid(); top_bar.fill.fore_color.rgb = col; top_bar.line.fill.background()

        txt(s, num,   x + Inches(0.12), Inches(2.3),  Inches(0.4),  Inches(0.4),
            size=Pt(18), bold=True, color=col)
        txt(s, title, x + Inches(0.12), Inches(2.72), Inches(2.7),  Inches(0.5),
            size=Pt(15), bold=True, color=WHITE)
        txt(s, body,  x + Inches(0.12), Inches(3.25), Inches(2.8),  Inches(2.7),
            size=Pt(12), color=GREY)

        # arrow between cards
        if i < 3:
            txt(s, "→", x + Inches(3.02), Inches(3.8), Inches(0.3), Inches(0.4),
                size=Pt(20), color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    txt(s, "Server calls POST /trigger → KeeperHub runs all 4 steps → pulls audit → mirrors to 0G Storage.",
        Inches(0.6), Inches(6.45), Inches(12), Inches(0.6),
        size=Pt(14), color=ACCENT2, italic=True, align=PP_ALIGN.CENTER)


def slide_stack(prs):
    """09 — Tech Stack"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "TECH STACK")

    txt(s, "What we built it with",
        Inches(0.6), Inches(0.85), Inches(12), Inches(0.65),
        size=Pt(28), bold=True, color=WHITE)

    rows = [
        ("Blockchain",    "0G Chain testnet · chainId 16602 · evmRPC + chainscan-galileo explorer"),
        ("Smart contracts","Solidity 0.8.24 · Hardhat 2 · OpenZeppelin v5 · evmVersion cancun"),
        ("Game server",   "Python 3.12 · FastAPI · uvicorn · pydantic · web3.py"),
        ("Backgammon engine","GNU Backgammon (gnubg) via External Player socket interface"),
        ("0G Storage bridge","Node.js workspace package (og-bridge) · @0gfoundation/0g-ts-sdk · ethers"),
        ("Identity",      "ENS subnames · text records (elo, match_count, style_uri, archive_uri)"),
        ("Settlement",    "KeeperHub workflows · kh CLI · HTTP trigger API"),
        ("Frontend",      "Next.js 16 · React 19 · TypeScript · Tailwind CSS 4 · wagmi 3 · viem 2"),
        ("Testing",       "82 server tests (pytest) · 52 contract tests (Hardhat/Mocha) · TDD throughout"),
    ]

    for i, (cat, detail) in enumerate(rows):
        y = Inches(1.75 + i * 0.6)
        txt(s, cat,    Inches(0.6), y, Inches(2.5),  Inches(0.55),
            size=Pt(13), bold=True, color=ACCENT2)
        txt(s, detail, Inches(3.2), y, Inches(9.8), Inches(0.55),
            size=Pt(13), color=WHITE)


def slide_status(prs):
    """10 — What's Built + Roadmap"""
    s = blank_slide(prs)
    bg(s); accent_bar(s)
    section_label(s, "STATUS & ROADMAP")

    txt(s, "Phases complete today",
        Inches(0.5), Inches(0.85), Inches(6.0), Inches(0.65),
        size=Pt(22), bold=True, color=WHITE)

    done = [
        "Phase 0 — Scaffold: pnpm workspace, uv server, Hardhat contracts, Next.js frontend",
        "Phase 1 — gnubg wrapper: FastAPI + pexpect; full match end-to-end",
        "Phase 2 — Core contracts: EloMath, MatchRegistry, AgentRegistry (52 tests)",
        "Phase 3 — gameRecordHash field on MatchRegistry",
        "Phase 4 — Deploy + verify on 0G testnet (chainscan-galileo)",
        "Phase 5 — AgentRegistry → ERC-7857 iNFT (tier + dataHashes)",
        "Phase 6 — 0G Storage round-trip via og-bridge Node helper",
        "Phase 7 — GameRecord upload + recordMatch with rootHash",
        "Phase 8 — Encrypted gnubg weights on 0G Storage (AES-256-GCM)",
        "Phase 9 — Agent experience overlay: learns after each match (82 tests)",
    ]
    bullet_frame(s, done, Inches(0.5), Inches(1.6), Inches(6.2), Inches(5.5),
                 size=Pt(11.5), color=ACCENT2, marker="✅ ")

    txt(s, "Coming next",
        Inches(7.0), Inches(0.85), Inches(5.8), Inches(0.65),
        size=Pt(22), bold=True, color=WHITE)

    next_up = [
        "Phase 10–11 — ENS subnames + text record updates",
        "Phase 12–15 — Frontend: wallet, agents list, match flow, ENS display",
        "Phase 16–19 — KeeperHub workflow (full settlement)",
        "Phase 20–21 — Match replay + audit trail display",
        "Phase 22 — KeeperHub feedback doc (bounty)",
        "Phase 24–25 — Deploy + demo video + submission",
    ]
    bullet_frame(s, next_up, Inches(7.0), Inches(1.6), Inches(5.8), Inches(3.5),
                 size=Pt(13), color=GREY, marker="⬜ ")

    # v2/v3 note
    v_box = s.shapes.add_shape(1, Inches(7.0), Inches(5.2), Inches(5.8), Inches(1.9))
    v_box.fill.solid(); v_box.fill.fore_color.rgb = RGBColor(0x14, 0x12, 0x28)
    v_box.line.color.rgb = ACCENT; v_box.line.width = Pt(1)
    txt(s, "v2: commit-reveal dice (VRF) · anti-cheat · L2 ENS\nv3: agent-vs-agent tournaments · zkML proofs · 0G Compute\nv4: open agent marketplace — bring your own engine",
        Inches(7.15), Inches(5.3), Inches(5.5), Inches(1.7),
        size=Pt(12), color=GREY, italic=True)


# ── main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    import os
    prs = new_prs()

    slide_title(prs)
    slide_problem(prs)
    slide_market(prs)
    slide_solution(prs)
    slide_how_it_works(prs)
    slide_ens(prs)
    slide_0g(prs)
    slide_keeperhub_intro(prs)
    slide_agents(prs)
    slide_elo(prs)
    slide_contracts(prs)
    slide_storage(prs)
    slide_keeperhub(prs)
    slide_stack(prs)
    slide_status(prs)

    out = os.path.join(os.path.dirname(__file__), "..", "chaingammon.pptx")
    prs.save(out)
    print(f"Saved: {os.path.abspath(out)}")
